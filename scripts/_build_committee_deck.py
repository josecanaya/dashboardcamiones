# -*- coding: utf-8 -*-
"""Construye los entregables del Comité (PPTX + MD + XLSX base) desde la salida
del ETL del repo (scripts/output/committee/) + Excel Movimientos.
Estilo aproximado a data/ejemplos/Comite 24-06 (2).pdf."""
import csv, json, os, io, glob
from collections import defaultdict
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import openpyxl

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CM = os.path.join(ROOT, "scripts", "output", "committee")
IMG = os.path.join(ROOT, "scripts", "output", "committee", "_charts")
DECK_IMG = r"C:\Users\Pc\AppData\Local\Temp\claude\C--Users-Pc-Desktop-JOSE-Programacion-Vicentin-rutas-dashboard\35ffd9cf-590b-4786-aa12-f352d6812b9a\scratchpad\deck_img"
OUT = os.path.join(ROOT, "salida_comite")
os.makedirs(IMG, exist_ok=True); os.makedirs(OUT, exist_ok=True)
PERIOD_A, PERIOD_B = "22/06/2026", "05/07/2026"
STEM = "Comite_Vicentin_2026-06-22_a_2026-07-05"

# ---------- Paleta ----------
NAVY=RGBColor(0x1F,0x3A,0x56); BLUE=RGBColor(0x2B,0x6C,0xB0); GREEN=RGBColor(0x2E,0x7D,0x32)
VIC=RGBColor(0x1C,0x6B,0x3C); ORANGE=RGBColor(0xDE,0x95,0x2E); GREY=RGBColor(0x64,0x74,0x8B)
LGREY=RGBColor(0x94,0xA3,0xB8); WHITE=RGBColor(0xFF,0xFF,0xFF); CARDBORDER=RGBColor(0xE2,0xE8,0xF0)
CARDBG=RGBColor(0xF8,0xFA,0xFC); RED=RGBColor(0xE4,0x57,0x4C)
SEG_COLORS=["#3B6FD4","#E4574C","#6FBF73","#9B7FD4","#F5B841","#4FB0A3","#7FB0E8"]

def hx(c): return "#%02X%02X%02X"%(c[0],c[1],c[2])

# ---------- Cargar datos ETL ----------
st=json.load(open(os.path.join(CM,"_stats.json"),encoding="utf-8"))
EX=st["executive"]; CO=st["coherence"]; MC=st["movimientosContrato"]; LIQ=MC.get("liquidMovements",{})

def read_csv(name):
    p=os.path.join(CM,name)
    if not os.path.exists(p): return []
    return list(csv.DictReader(open(p,encoding="utf-8")))

circ_sum={r["executive_circuit_code"]:r for r in read_csv("circuit_timing_summary.csv")}
seg_kpi=read_csv("segment_timing_kpi_COMPLETOS.csv")
legs=read_csv("segment_timing_legs_COMPLETOS.csv")

def fnum(x):
    try: return float(x)
    except: return 0.0

# tramos por circuito
seg_by_circ=defaultdict(dict)
for r in seg_kpi:
    seg_by_circ[r["executive_circuit_code"]][(r["from_logical"],r["to_logical"])]=r
legs_by_circ_trans=defaultdict(list)
for r in legs:
    legs_by_circ_trans[(r["executive_circuit_code"],r["from_logical"],r["to_logical"])].append(fnum(r["duration_min"]))

# ---------- Volúmenes desde Excel (tonelaje) ----------
import openpyxl as ox
MOVDIR=os.path.join(ROOT,"data","Movimientos")
prod_kg=defaultdict(float); prod_n=defaultdict(int)
planta_n=defaultdict(int); mov_n=defaultdict(int); aceite_kg=0.0; aceite_n=0
LIQ_RE=("ACEITE","GLICERINA","BORRA","LECITINA","NITROGENO LIQUIDO","SODA CAUSTICA","ALCOHOL","FUEL","GAS OIL","HIPOCLORITO")
def is_liquid(p): return any(k in p for k in LIQ_RE)
mov_rows=0
for f in sorted(glob.glob(os.path.join(MOVDIR,"MovimientosPorContrato_*.xlsx"))):
    wb=ox.load_workbook(f,read_only=True,data_only=True); ws=wb.active
    rows=ws.iter_rows(values_only=True); header=next(rows)
    idx={str(h).strip():i for i,h in enumerate(header)}
    def g(row,k):
        i=idx.get(k); return row[i] if i is not None and i<len(row) else None
    for row in rows:
        if row is None or all(c is None for c in row): continue
        prod=str(g(row,"Producto") or "").strip()
        if not prod: continue
        mov_rows+=1
        kg=fnum(g(row,"Kgs.Neto"))
        prod_kg[prod]+=kg; prod_n[prod]+=1
        planta_n[str(g(row,"Planta") or "").strip()]+=1
        mov_n[str(g(row,"Mov") or "").strip()]+=1
        if is_liquid(prod):
            global_a=True
    wb.close()
# recompute aceite
for p,kg in prod_kg.items():
    if is_liquid(p): aceite_kg+=kg; aceite_n+=prod_n[p]
tot_kg=sum(prod_kg.values()); tot_mov=sum(prod_n.values())
soja_kg=prod_kg.get("SOJA",0); gira_kg=prod_kg.get("GIRASOL",0)+prod_kg.get("ACEITE GIRASOL PARA REFINAR",0)
solidos_n=sum(n for p,n in prod_n.items() if not is_liquid(p))
liquidos_n=sum(n for p,n in prod_n.items() if is_liquid(p))

print(f"tot_mov={tot_mov} tot_kt={tot_kg/1e6:.1f} aceite_n={aceite_n} aceite_kt={aceite_kg/1e6:.1f} solidos={solidos_n} liquidos={liquidos_n}")

# ================= CHARTS (matplotlib) =================
plt.rcParams.update({"font.family":"DejaVu Sans","font.size":11})

def save(fig,name):
    p=os.path.join(IMG,name); fig.savefig(p,dpi=170,bbox_inches="tight",transparent=True,pad_inches=0.05); plt.close(fig); return p

def chart_stacked_segments(steps, name):
    """steps: list of (label, minutes, color). Barra horizontal apilada por % del total."""
    total=sum(m for _,m,_ in steps) or 1
    fig,ax=plt.subplots(figsize=(11,1.5)); left=0
    for lab,m,c in steps:
        pct=m/total*100
        ax.barh(0,pct,left=left,color=c,edgecolor="white",height=0.7)
        if pct>=4:
            ax.text(left+pct/2,0,f"{pct:.1f}%",ha="center",va="center",color="white",fontsize=11,fontweight="bold")
        left+=pct
    ax.set_xlim(0,100); ax.set_ylim(-0.5,0.5); ax.set_yticks([])
    ax.set_xticks([0,20,40,60,80,100]); ax.set_xticklabels(["0%","20%","40%","60%","80%","100%"],fontsize=10,color=hx(GREY))
    for s in ["top","right","left"]: ax.spines[s].set_visible(False)
    ax.spines["bottom"].set_color("#CBD5E1")
    return save(fig,name)

def chart_scatter(durations, mean, name, color="#3B6FD4"):
    """Scatter estilo deck: dispersión de duraciones (min) con curva normal."""
    d=np.array([x for x in durations if x>0])
    if len(d)==0: d=np.array([mean])
    rng=np.random.default_rng(7); y=rng.uniform(0,10,size=len(d))
    fig,ax=plt.subplots(figsize=(6.2,3.4))
    ax.scatter(d,y,s=14,color=color,alpha=0.55,edgecolors="none")
    xs=np.linspace(0,max(d.max(),mean*2),200); mu=d.mean(); sd=d.std() or 1
    curve=np.exp(-0.5*((xs-mu)/sd)**2); curve=curve/curve.max()*10
    ax.plot(xs,curve,color="#7FA8D8",lw=1.6)
    ax.axvline(mu,color="#2B6CB0",lw=1.4)
    ax.set_ylim(0,11); ax.set_xlim(0,xs.max()); ax.set_yticks([])
    ax.set_xlabel("Minutos",fontsize=9,color=hx(GREY)); ax.set_ylabel("Camiones",fontsize=9,color=hx(GREY))
    for s in ["top","right"]: ax.spines[s].set_visible(False)
    for s in ["bottom","left"]: ax.spines[s].set_color("#CBD5E1")
    ax.tick_params(colors=hx(GREY),labelsize=8)
    return save(fig,name)

def chart_donut(values,labels,colors,name):
    fig,ax=plt.subplots(figsize=(3.2,3.2))
    vals=[max(v,0.0001) for v in values]
    w,_=ax.pie(vals,colors=colors,startangle=90,counterclock=False,wedgeprops=dict(width=0.42,edgecolor="white"))
    ax.set(aspect="equal")
    return save(fig,name)

def chart_hbars(items,name,color="#2E7D32",unit="%"):
    """items: list of (label, pct, rightlabel). Etiqueta a la izquierda del track."""
    fig,ax=plt.subplots(figsize=(6.6,0.62*len(items)+0.25))
    ys=range(len(items))[::-1]
    for y,(lab,pct,rl) in zip(ys,items):
        ax.barh(y,100,color="#EEF2F6",height=0.55,zorder=1)
        ax.barh(y,pct,color=color,height=0.55,zorder=2)
        ax.text(-4,y,f"{lab}",va="center",ha="right",fontsize=10,color=hx(NAVY),fontweight="bold",zorder=4)
        if pct>=16:
            ax.text(pct-2,y,f"{pct:.0f}%",va="center",ha="right",fontsize=9,color="white",fontweight="bold",zorder=4)
        else:
            ax.text(pct+2,y,f"{pct:.0f}%" if pct>=10 else f"{pct:.1f}%",va="center",ha="left",fontsize=9,color=hx(GREY),fontweight="bold",zorder=4)
        ax.text(104,y,f"{rl}",va="center",ha="left",fontsize=9,color=hx(GREY))
    ax.set_xlim(-62,118); ax.set_ylim(-0.6,len(items)-0.4); ax.axis("off")
    return save(fig,name)

print("charts helpers ready")

# ============ PPTX ============
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
SW,SH=prs.slide_width,prs.slide_height

def slide(): return prs.slides.add_slide(BLANK)

def _set(fr,text,size,color,bold=False,align=PP_ALIGN.LEFT,font="Calibri"):
    p=fr.paragraphs[0]; p.alignment=align
    r=p.runs[0] if p.runs else p.add_run()
    r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; r.font.name=font

def textbox(sl,x,y,w,h,text,size,color,bold=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,font="Calibri"):
    tb=sl.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    _set(tf,text,size,color,bold,align,font); return tb

def add_multiline(tb,lines):
    tf=tb.text_frame
    for i,(text,size,color,bold) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        r=p.runs[0] if p.runs else p.add_run()
        r.text=text; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; r.font.name="Calibri"

def rrect(sl,x,y,w,h,fill,line=None,line_w=1.0,radius=0.08):
    sp=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,x,y,w,h)
    sp.adjustments[0]=radius
    sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(line_w)
    sp.shadow.inherit=False
    return sp

def rect(sl,x,y,w,h,fill):
    sp=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,x,y,w,h)
    sp.fill.solid(); sp.fill.fore_color.rgb=fill; sp.line.fill.background(); sp.shadow.inherit=False
    return sp

def underline(sl,x,y,w):
    # gradiente azul->verde simulado con dos rectángulos
    rect(sl,x,y,int(w*0.5),Pt(4),BLUE); rect(sl,x+int(w*0.5),y,int(w*0.5),Pt(4),GREEN)

def title_bar(sl,text):
    textbox(sl,Inches(1.0),Inches(0.28),Inches(11.3),Inches(0.7),text,26,NAVY,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
    underline(sl,Inches(4.9),Inches(1.02),Inches(3.5))

def logo(sl):
    lg=os.path.join(ROOT,"public","logo.png")
    # marca textual esquina
    tb=rrect(sl,Inches(0.35),Inches(0.28),Inches(1.25),Inches(0.62),VIC,radius=0.12)
    textbox(sl,Inches(0.4),Inches(0.31),Inches(1.15),Inches(0.56),"Nueva\nVicentin\nArgentina.",8,WHITE,True,PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)

def pic(sl,path,x,y,w=None,h=None):
    if not os.path.exists(path): return None
    return sl.shapes.add_picture(path,x,y,width=w,height=h)

# ---- diccionario de números ----
def R(v,dec=0):
    return f"{v:,.{dec}f}".replace(",","·").replace(".",",").replace("·",".")

print("pptx helpers ready")

# ============ DATA MODEL DECK ============
def circ_total(code):
    r=circ_sum.get(code); return (fnum(r["mean_total_min"]),fnum(r["std_total_min"]),int(fnum(r["n_journeys"]))) if r else (0,0,0)

R7_STEPS=[
 ("Ingreso","Ingreso → Preingreso",("INGRESO","PREINGRESO"),SEG_COLORS[0]),
 ("Calada","Preingreso → Calada",("PREINGRESO","CALADA"),SEG_COLORS[1]),
 ("Egreso","Calada → Egreso",("CALADA","EGRESO"),SEG_COLORS[2]),
 ("Tránsito interpl.","Egreso → Ingreso SL",("EGRESO","SL_INGRESO"),SEG_COLORS[3]),
 ("Espera Balanza","Ingreso → Balanza",("SL_INGRESO","SL_BALANZA_INGRESO"),SEG_COLORS[4]),
 ("Descarga","Balanza In → Balanza Eg",("SL_BALANZA_INGRESO","SL_EGRESO"),SEG_COLORS[5]),
]
def step_val(code,trans):
    r=seg_by_circ.get(code,{}).get(trans)
    if not r: return (0,0,0)
    return (fnum(r["mean_min"]),fnum(r["std_min"]),int(fnum(r["n"])))

# ---- Sample numbers ----
soja_real=prod_n.get("SOJA",0); soja_capt=circ_total("R7")[2]
gira_real=prod_n.get("GIRASOL",0); r5n=circ_total("R5")[2]; r6n=circ_total("R6")[2]; gira_capt=r5n+r6n
r1n=circ_total("R1")[2]
total_capt=soja_capt+gira_capt; total_real=soja_real+gira_real
def pct(a,b): return (a/b*100) if b else 0
soja_kt=soja_kg/1e6; gira_kt=prod_kg.get("GIRASOL",0)/1e6

# =================================================================
# SLIDE 1 — Portada
# =================================================================
s=slide()
bg=os.path.join(DECK_IMG,"past_02.png")
if not os.path.exists(bg): bg=os.path.join(DECK_IMG,"past_01.png")
if os.path.exists(bg): s.shapes.add_picture(bg,0,0,width=SW,height=SH)
ov=rect(s,0,0,SW,SH,RGBColor(0x0F,0x1F,0x2E)); ov.fill.fore_color.rgb=RGBColor(0x0F,0x1F,0x2E)
ov.fill.transparency=0  # pptx no soporta transp directa; usar caja translúcida via alpha en xml
# pill superior
rrect(s,Inches(0.5),Inches(0.55),Inches(3.5),Inches(0.5),WHITE,radius=0.25)
textbox(s,Inches(0.6),Inches(0.58),Inches(3.4),Inches(0.44),"Comité de seguridad y eficiencia",12,NAVY,False,PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)
textbox(s,Inches(1.0),Inches(2.7),Inches(11.3),Inches(1.2),"Nueva Vicentin Argentina",44,WHITE,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
underline(s,Inches(5.2),Inches(3.95),Inches(2.9))
rrect(s,Inches(0.5),Inches(6.5),Inches(2.4),Inches(0.5),WHITE,radius=0.2)
textbox(s,Inches(0.6),Inches(6.53),Inches(2.3),Inches(0.44),f"{PERIOD_A}  —  {PERIOD_B}",12,NAVY,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
rrect(s,Inches(11.1),Inches(6.45),Inches(1.75),Inches(0.7),VIC,radius=0.12)
textbox(s,Inches(11.2),Inches(6.5),Inches(1.6),Inches(0.6),"Nueva\nVicentin\nArgentina.",10,WHITE,True,PP_ALIGN.LEFT,MSO_ANCHOR.MIDDLE)

# =================================================================
# SLIDE 2 — Tamaño de la Muestra y Periodo
# =================================================================
s=slide(); logo(s); title_bar(s,"Tamaño de la Muestra y Periodo")
def stat_card(x,color,titulo,valor,sub1,sub2):
    rrect(s,x,Inches(1.5),Inches(3.6),Inches(1.65),CARDBG,CARDBORDER,1.2,0.08)
    textbox(s,x,Inches(1.62),Inches(3.6),Inches(0.35),titulo,12,color,True,PP_ALIGN.CENTER)
    textbox(s,x,Inches(1.95),Inches(3.6),Inches(0.65),valor,34,color,True,PP_ALIGN.CENTER)
    textbox(s,x,Inches(2.6),Inches(3.6),Inches(0.3),sub1,11,GREY,False,PP_ALIGN.CENTER)
    textbox(s,x,Inches(2.82),Inches(3.6),Inches(0.3),sub2,10,GREY,False,PP_ALIGN.CENTER)
stat_card(Inches(0.9),BLUE,"TOTAL CAPTADO",R(total_capt),"camiones",f"Total real: {R(total_real)} ({pct(total_capt,total_real):.0f}%)")
stat_card(Inches(4.85),GREEN,"SOJA CAPTADA",R(soja_capt),"camiones",f"Total real: {R(soja_real)} ({pct(soja_capt,soja_real):.0f}%)")
stat_card(Inches(8.8),ORANGE,"GIRASOL CAPTADO",R(gira_capt),"camiones",f"Total real: {R(gira_real)} ({pct(gira_capt,gira_real):.0f}%)")
# panel izq: disponibilidad
textbox(s,Inches(0.9),Inches(3.4),Inches(6.5),Inches(0.35),"DISPONIBILIDAD CAPTADA COMO % DEL TOTAL REAL",12,NAVY,True)
textbox(s,Inches(0.9),Inches(3.85),Inches(6.5),Inches(0.3),f"SOJA — {R(soja_capt)} captados de {R(soja_real)} reales ({pct(soja_capt,soja_real):.1f}%)",11,GREEN,True)
soja_items=[("R7 · San Lorenzo",pct(soja_capt,soja_capt+r1n),R(soja_capt)),("R1 · Celda 16",pct(r1n,soja_capt+r1n),R(r1n))]
p=chart_hbars(soja_items,"soja_avail.png",color=hx(GREEN)); s.shapes.add_picture(p,Inches(1.1),Inches(4.2),width=Inches(5.7))
textbox(s,Inches(0.9),Inches(5.55),Inches(6.5),Inches(0.3),f"GIRASOL — {R(gira_capt)} captados de {R(gira_real)} reales ({pct(gira_capt,gira_real):.1f}%)",11,ORANGE,True)
gira_items=[("R5 · Volcable 1",pct(r5n,gira_capt) if gira_capt else 0,R(r5n)),("R6 · Volcable 2",pct(r6n,gira_capt) if gira_capt else 0,R(r6n))]
p=chart_hbars(gira_items,"gira_avail.png",color=hx(ORANGE)); s.shapes.add_picture(p,Inches(1.1),Inches(5.9),width=Inches(5.7))
# panel der: muestra seleccionada
rrect(s,Inches(7.7),Inches(3.4),Inches(4.9),Inches(3.4),CARDBG,CARDBORDER,1.2,0.05)
textbox(s,Inches(7.7),Inches(3.55),Inches(4.9),Inches(0.35),"MUESTRA SELECCIONADA PARA EL ANÁLISIS",12,NAVY,True,PP_ALIGN.CENTER)
textbox(s,Inches(7.7),Inches(3.85),Inches(4.9),Inches(0.3),"Se analiza la muestra clasificada del período",9,GREY,False,PP_ALIGN.CENTER)
rrect(s,Inches(8.0),Inches(4.35),Inches(4.3),Inches(1.05),WHITE,CARDBORDER,1,0.08)
textbox(s,Inches(8.2),Inches(4.5),Inches(2.9),Inches(0.4),"SOJA · R7",14,GREEN,True)
textbox(s,Inches(8.2),Inches(4.9),Inches(2.9),Inches(0.35),f"{pct(soja_capt,soja_real):.0f}% del real de R7",10,GREY)
textbox(s,Inches(10.6),Inches(4.45),Inches(1.6),Inches(0.7),R(soja_kt,1),24,GREEN,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
textbox(s,Inches(10.6),Inches(5.05),Inches(1.6),Inches(0.3),"Kilotoneladas",9,GREY,False,PP_ALIGN.CENTER)
rrect(s,Inches(8.0),Inches(5.55),Inches(4.3),Inches(1.05),WHITE,CARDBORDER,1,0.08)
textbox(s,Inches(8.2),Inches(5.7),Inches(2.9),Inches(0.4),"GIRASOL · R5/R6",14,ORANGE,True)
textbox(s,Inches(8.2),Inches(6.1),Inches(2.9),Inches(0.35),f"{pct(gira_capt,gira_real):.0f}% del real de R5+R6",10,GREY)
textbox(s,Inches(10.6),Inches(5.65),Inches(1.6),Inches(0.7),R(gira_kt,1),24,ORANGE,True,PP_ALIGN.CENTER,MSO_ANCHOR.MIDDLE)
textbox(s,Inches(10.6),Inches(6.25),Inches(1.6),Inches(0.3),"Kilotoneladas",9,GREY,False,PP_ALIGN.CENTER)

print("slides 1-2 ok")
