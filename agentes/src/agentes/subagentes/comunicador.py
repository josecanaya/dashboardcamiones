"""Comunicador: PPTX de comité a partir de get_summary (sin reglas de negocio)."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from agentes.etl_client import EtlClient


def generar_pptx_comite(
    run_id: str,
    *,
    client: EtlClient | None = None,
    out_dir: str | Path | None = None,
) -> dict[str, Any]:
    """Genera un PPTX mínimo con KPIs del summary de la corrida."""
    c = client or EtlClient()
    summary = c.get_summary(run_id)
    stats = summary.get("stats") or {}
    ex = stats.get("executive") or {}

    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = f"Comité ETL — {run_id}"
    body = slide.placeholders[1].text_frame
    body.clear()

    lines = [
        f"Reglas: {(summary.get('manifest') or {}).get('rulesVersion') or 'n/d'}",
        f"Completos (comité): {ex.get('committeeCompletos', 'n/d')}",
        f"Variaciones: {ex.get('committeeVariaciones', 'n/d')}",
        f"Anomalías (comité): {ex.get('committeeAnomalias', 'n/d')}",
        f"Incompletos: {ex.get('incompletos', 'n/d')}",
        f"Deducidos: {ex.get('deducidos', 'n/d')}",
        f"Eventos: {ex.get('eventCount', 'n/d')}",
    ]
    for i, line in enumerate(lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)

    dest_dir = Path(out_dir or Path.cwd() / "output" / "comite")
    dest_dir.mkdir(parents=True, exist_ok=True)
    dest = dest_dir / f"comite_{run_id}.pptx"
    prs.save(str(dest))
    return {"run_id": run_id, "path": str(dest.resolve()), "kpis": ex}


PPTX_TOOL = {
    "name": "generar_pptx_comite",
    "description": (
        "Genera una presentación PPTX de comité con KPIs de get_summary para un run_id. "
        "Usar cuando pidan slides o material para el comité."
    ),
    "input_schema": {
        "type": "object",
        "properties": {
            "run_id": {"type": "string"},
            "out_dir": {"type": "string", "description": "Carpeta de salida (opcional)."},
        },
        "required": ["run_id"],
        "additionalProperties": False,
    },
}
